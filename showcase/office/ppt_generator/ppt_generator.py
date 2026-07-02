"""
[Showcase / Office / 通用]
Skill: 从结构化内容生成 PPT

适用场景：
  周报/月报这类固定格式的汇报 PPT，内容每次变、排版不变，
  直接用 python-pptx 生成，不需要打开 PowerPoint / Keynote，不需要屏幕。

用法：
  python run.py showcase/office/ppt_generator/ppt_generator -- \\
    --output weekly_report.pptx \\
    --data '[{"title":"本周进展","bullets":["完成 A","完成 B"]},{"title":"下周计划","bullets":["开始 C"]}]'
"""

import argparse
import json
import sys

from pptx import Presentation

from core.logger import SkillLogger


def generate_ppt(slides: list[dict], output: str) -> None:
    prs = Presentation()
    bullet_layout = prs.slide_layouts[1]  # 标题 + 内容布局

    for slide_data in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        if bullets:
            body = slide.placeholders[1].text_frame
            body.text = bullets[0]
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet

    prs.save(output)


def main():
    parser = argparse.ArgumentParser(description="从结构化内容生成 PPT")
    parser.add_argument("--data", required=True, help='幻灯片内容 JSON 数组，每项 {"title":..., "bullets":[...]}')
    parser.add_argument("--output", required=True, help="输出的 .pptx 路径")
    try:
        sep = sys.argv.index("--")
        argv = sys.argv[sep + 1:]
    except ValueError:
        argv = sys.argv[1:]
    args = parser.parse_args(argv)

    log = SkillLogger("office/ppt_generator")

    slides = json.loads(args.data)
    log.step(f"生成 {len(slides)} 页 PPT")
    generate_ppt(slides, args.output)
    print(f"已保存到 {args.output}")
    log.finish({"slides": len(slides), "output": args.output})
    return {"output": args.output, "slides": len(slides)}


if __name__ == "__main__":
    main()
